import streamlit as st
import pandas as pd

from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer
)

from reportlab.lib.styles import (
    getSampleStyleSheet
)

from pptx import Presentation

# ---------------------------------------------------
# PAGE CONFIG
# ---------------------------------------------------

st.set_page_config(
    page_title="Export Reports",
    layout="wide"
)

# ---------------------------------------------------
# TITLE
# ---------------------------------------------------

st.title("📄 Report Export Engine")

st.markdown("""
Generate investor-style intelligence reports.
""")

# ---------------------------------------------------
# LOAD DATA
# ---------------------------------------------------

try:

    df = pd.read_csv(
        "data/historical_reviews.csv"
    )

except Exception as e:

    st.error(f"Data Load Error: {e}")

    st.stop()

# ---------------------------------------------------
# SUMMARY
# ---------------------------------------------------

summary = (
    df.groupby("platform")
    .agg({
        "rating": "mean",
        "review": "count"
    })
    .reset_index()
)

summary.columns = [
    "Platform",
    "Avg Rating",
    "Total Reviews"
]

# ---------------------------------------------------
# DISPLAY
# ---------------------------------------------------

st.subheader("📊 Current Intelligence Snapshot")

st.dataframe(
    summary,
    use_container_width=True
)

# ---------------------------------------------------
# PDF EXPORT
# ---------------------------------------------------

if st.button("📄 Generate PDF Report"):

    pdf_file = "wealthtech_report.pdf"

    doc = SimpleDocTemplate(pdf_file)

    styles = getSampleStyleSheet()

    elements = []

    title = Paragraph(
        "WealthTech Intelligence Report",
        styles["Title"]
    )

    elements.append(title)

    elements.append(Spacer(1, 12))

    for _, row in summary.iterrows():

        text = f"""
        <b>{row['Platform']}</b><br/>
        Avg Rating: {round(row['Avg Rating'],2)}<br/>
        Total Reviews: {row['Total Reviews']}<br/><br/>
        """

        elements.append(
            Paragraph(
                text,
                styles["BodyText"]
            )
        )

    doc.build(elements)

    st.success("PDF report generated.")

    with open(pdf_file, "rb") as f:

        st.download_button(
            "⬇ Download PDF",
            f,
            file_name=pdf_file
        )

# ---------------------------------------------------
# PPT EXPORT
# ---------------------------------------------------

if st.button("📊 Generate PowerPoint"):

    ppt_file = "wealthtech_report.pptx"

    prs = Presentation()

    slide_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(
        slide_layout
    )

    title = slide.shapes.title

    subtitle = slide.placeholders[1]

    title.text = (
        "WealthTech Intelligence Report"
    )

    subtitle.text = (
        "Advisor & Fintech Analytics"
    )

    # PLATFORM SLIDES
    for _, row in summary.iterrows():

        slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(
            slide_layout
        )

        title = slide.shapes.title

        content = slide.placeholders[1]

        title.text = row["Platform"]

        content.text = f"""
Average Rating:
{round(row['Avg Rating'],2)}

Total Reviews:
{row['Total Reviews']}
"""

    prs.save(ppt_file)

    st.success(
        "PowerPoint generated."
    )

    with open(ppt_file, "rb") as f:

        st.download_button(
            "⬇ Download PowerPoint",
            f,
            file_name=ppt_file
        )

# ---------------------------------------------------
# CSV EXPORT
# ---------------------------------------------------

csv = df.to_csv(index=False)

st.download_button(
    "⬇ Download Historical CSV",
    csv,
    file_name="historical_reviews.csv",
    mime="text/csv"
)
