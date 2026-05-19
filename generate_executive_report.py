import os
import pandas as pd

from datetime import datetime

from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer
)

from reportlab.lib.styles import (
    getSampleStyleSheet
)

from pptx import Presentation

# ---------------------------------------------------
# CREATE REPORT FOLDER
# ---------------------------------------------------

os.makedirs(
    "reports",
    exist_ok=True
)

# ---------------------------------------------------
# LOAD DATA
# ---------------------------------------------------

df = pd.read_csv(
    "data/historical_reviews.csv"
)

# ---------------------------------------------------
# SUMMARY
# ---------------------------------------------------

summary = (

    df.groupby("platform")

    .agg({

        "rating": "mean",

        "review": "count",

        "likes": "sum"

    })

    .reset_index()

)

summary.columns = [

    "Platform",

    "Average Rating",

    "Total Reviews",

    "Total Likes"

]

# ---------------------------------------------------
# NORMALIZED SCORE
# ---------------------------------------------------

summary["Review Score"] = (

    summary["Total Reviews"]

    / summary["Total Reviews"].max()

) * 100

summary["Rating Score"] = (

    summary["Average Rating"]

    / 5

) * 100

summary["Engagement Score"] = (

    summary["Total Likes"]

    / summary["Total Likes"].max()

) * 100

summary["Benchmark Score"] = (

    summary["Review Score"] * 0.4

    +

    summary["Rating Score"] * 0.4

    +

    summary["Engagement Score"] * 0.2

)

summary["Benchmark Score"] = (

    summary["Benchmark Score"]

    .round(2)

)

summary = summary.sort_values(

    by="Benchmark Score",

    ascending=False

)

# ---------------------------------------------------
# PDF REPORT
# ---------------------------------------------------

today = datetime.now().strftime("%Y_%m")

pdf_name = (
    f"reports/Executive_Report_{today}.pdf"
)

doc = SimpleDocTemplate(pdf_name)

styles = getSampleStyleSheet()

elements = []

# TITLE

title = Paragraph(

    "WealthTech Executive Intelligence Report",

    styles["Title"]

)

elements.append(title)

elements.append(
    Spacer(1, 20)
)

# MARKET SUMMARY

leader = summary.iloc[0]

market_text = f"""
<b>Market Overview</b><br/><br/>

Top Platform:
{leader['Platform']}<br/><br/>

Benchmark Score:
{leader['Benchmark Score']}<br/><br/>

Total Market Reviews:
{len(df)}<br/><br/>

The WealthTech ecosystem is rapidly evolving toward:
- advisor operating systems
- AI-powered workflows
- mobile-first onboarding
- automation-led engagement
"""

elements.append(

    Paragraph(
        market_text,
        styles["BodyText"]
    )

)

elements.append(
    Spacer(1, 20)
)

# PLATFORM DETAILS

for _, row in summary.iterrows():

    text = f"""
<b>{row['Platform']}</b><br/><br/>

Average Rating:
{round(row['Average Rating'],2)}<br/>

Total Reviews:
{row['Total Reviews']}<br/>

Benchmark Score:
{row['Benchmark Score']}<br/><br/>
"""

    elements.append(

        Paragraph(
            text,
            styles["BodyText"]
        )

    )

    elements.append(
        Spacer(1, 12)
    )

# BUILD PDF

doc.build(elements)

# ---------------------------------------------------
# POWERPOINT REPORT
# ---------------------------------------------------

prs = Presentation()

# TITLE SLIDE

slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(
    slide_layout
)

title = slide.shapes.title

subtitle = slide.placeholders[1]

title.text = (
    "WealthTech Executive Intelligence"
)

subtitle.text = (
    f"Automated Benchmark Report - {today}"
)

# SUMMARY SLIDE

slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(
    slide_layout
)

title = slide.shapes.title

content = slide.placeholders[1]

title.text = "Market Overview"

content.text = f"""
Top Platform:
{leader['Platform']}

Benchmark Score:
{leader['Benchmark Score']}

Total Reviews:
{len(df)}

Leading platforms continue to gain momentum through:
- superior UX
- onboarding efficiency
- advisor engagement
- operational scalability
"""

# PLATFORM SLIDES

for _, row in summary.iterrows():

    slide = prs.slides.add_slide(
        slide_layout
    )

    title = slide.shapes.title

    content = slide.placeholders[1]

    title.text = row["Platform"]

    content.text = f"""
Average Rating:
{round(row['Average Rating'],2)}

Total Reviews:
{row['Total Reviews']}

Benchmark Score:
{row['Benchmark Score']}

AI Commentary:
Strong WealthTech positioning with growing advisor engagement.
"""

# SAVE PPT

ppt_name = (
    f"reports/Executive_Report_{today}.pptx"
)

prs.save(ppt_name)

# ---------------------------------------------------
# SUCCESS
# ---------------------------------------------------

print(
    f"PDF generated: {pdf_name}"
)

print(
    f"PPT generated: {ppt_name}"
)
